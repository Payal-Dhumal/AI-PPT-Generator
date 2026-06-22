"""PowerPoint generation engine."""

import os
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from app.models.presentation import PresentationResponse
from app.models.slide import SlideData
from app.services.image_service import ImageManager
from app.templates.base import TemplateConfig
from app.templates import get_template_config
from app.utils.constants import (
    BULLET_INDENT,
    CONTENT_MARGIN,
    GENERATED_DIR,
    IMAGE_MARGIN,
    IMAGE_PADDING,
    SLIDE_PADDING,
    TITLE_MARGIN,
)
from app.utils.helpers import sanitize_filename
from app.utils.validators import ValidationError


class PPTGenerator:
    """Creates export-ready PPTX files from structured slide data."""

    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    def __init__(
        self,
        image_manager: ImageManager,
        generated_dir: str = GENERATED_DIR,
    ):
        self.image_manager = image_manager
        self.generated_dir = generated_dir
        os.makedirs(self.generated_dir, exist_ok=True)

    def generate_presentation(
        self,
        presentation_data: PresentationResponse,
        template_name: str,
    ) -> str:
        template = get_template_config(template_name)
        prs = Presentation()
        prs.slide_width = self.SLIDE_WIDTH
        prs.slide_height = self.SLIDE_HEIGHT

        for slide_data in presentation_data.slides:
            self.create_slide(prs, slide_data, template)

        filename = sanitize_filename(presentation_data.presentation_title) + ".pptx"
        file_path = os.path.join(self.generated_dir, filename)

        counter = 1
        base_path = file_path
        while os.path.exists(file_path):
            name, ext = os.path.splitext(base_path)
            file_path = f"{name}_{counter}{ext}"
            counter += 1

        try:
            prs.save(file_path)
        except Exception as exc:
            raise ValidationError(f"PPT export failed: {exc}") from exc

        return file_path

    def create_slide(
        self,
        prs: Presentation,
        slide_data: SlideData,
        template: TemplateConfig,
    ) -> None:
        slide_type = slide_data.slide_type.lower()
        handlers = {
            "title": self._create_title_slide,
            "content": self._create_content_slide,
            "two_column": self._create_two_column_slide,
            "image_focus": self._create_image_focus_slide,
            "comparison": self._create_comparison_slide,
            "timeline": self._create_timeline_slide,
            "conclusion": self._create_conclusion_slide,
            "thank_you": self._create_thank_you_slide,
        }
        handler = handlers.get(slide_type, self._create_content_slide)
        handler(prs, slide_data, template)

    def _add_blank_slide(self, prs: Presentation):
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        self._apply_background(slide, prs)
        return slide

    def _apply_background(self, slide, prs: Presentation) -> None:
        background = slide.background
        fill = background.fill
        fill.solid()

    def _set_slide_background(self, slide, color_hex: str) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self._hex_to_rgb(color_hex)

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        hex_color = hex_color.lstrip("#")
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16),
        )

    def _add_text_box(
        self,
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        text: str,
        font_name: str,
        font_size: int,
        color_hex: str,
        bold: bool = False,
        alignment: PP_ALIGN = PP_ALIGN.LEFT,
    ):
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = self._hex_to_rgb(color_hex)
        p.alignment = alignment
        return textbox

    def _add_bullet_points(
        self,
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        items: List[str],
        template: TemplateConfig,
    ):
        textbox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = textbox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.name = template.body_font
            p.font.size = Pt(template.content_size)
            p.font.color.rgb = self._hex_to_rgb(template.content_color)
            p.space_after = Pt(8)

    def _insert_image(
        self,
        slide,
        image_name: str,
        position: str,
        content_bounds: Optional[Tuple[float, float, float, float]] = None,
    ) -> None:
        try:
            image_path = self.image_manager.get_image_path(image_name)
        except ValidationError:
            return

        slide_w = 13.333
        slide_h = 7.5
        margin = IMAGE_MARGIN

        position_map = {
            "left": (margin, margin + 1.2, 4.5, 4.5),
            "right": (slide_w - 4.5 - margin, margin + 1.2, 4.5, 4.5),
            "center": ((slide_w - 5) / 2, (slide_h - 4) / 2, 5, 4),
            "top": (margin + 1, margin + 0.8, slide_w - 2 * margin - 1, 3),
            "bottom": (margin + 1, slide_h - 3.5, slide_w - 2 * margin - 1, 2.5),
            "background": (0, 0, slide_w, slide_h),
        }

        if content_bounds and position in ("left", "right"):
            left, top, width, height = content_bounds
            if position == "right":
                coords = (left + width + IMAGE_PADDING, top, 4.0, height)
            else:
                coords = (margin, top, 4.0, height)
        else:
            coords = position_map.get(position, position_map["right"])

        left, top, width, height = coords
        slide.shapes.add_picture(image_path, Inches(left), Inches(top), Inches(width), Inches(height))

    def _create_title_slide(
        self, prs: Presentation, slide_data: SlideData, template: TemplateConfig
    ) -> None:
        slide = self._add_blank_slide(prs)
        self._set_slide_background(slide, template.background_color)

        self._add_text_box(
            slide,
            TITLE_MARGIN,
            2.0,
            13.333 - 2 * TITLE_MARGIN,
            1.5,
            slide_data.title,
            template.title_font,
            template.title_size,
            template.title_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        if slide_data.subtitle:
            self._add_text_box(
                slide,
                TITLE_MARGIN,
                3.8,
                13.333 - 2 * TITLE_MARGIN,
                1.0,
                slide_data.subtitle,
                template.body_font,
                template.subtitle_size,
                template.accent_color,
                alignment=PP_ALIGN.CENTER,
            )

        for img_ref in slide_data.images:
            self._insert_image(slide, img_ref.image_name, img_ref.position)

    def _create_content_slide(
        self, prs: Presentation, slide_data: SlideData, template: TemplateConfig
    ) -> None:
        slide = self._add_blank_slide(prs)
        self._set_slide_background(slide, template.background_color)

        has_right_image = any(img.position == "right" for img in slide_data.images)
        has_left_image = any(img.position == "left" for img in slide_data.images)
        content_width = 7.5 if (has_right_image or has_left_image) else 11.5
        content_left = 5.5 if has_left_image else TITLE_MARGIN

        self._add_text_box(
            slide,
            TITLE_MARGIN,
            CONTENT_MARGIN,
            11.5,
            0.8,
            slide_data.title,
            template.heading_font,
            template.heading_size,
            template.heading_color,
            bold=True,
        )

        content_top = CONTENT_MARGIN + 1.0
        self._add_bullet_points(
            slide,
            content_left,
            content_top,
            content_width,
            5.0,
            slide_data.content,
            template,
        )

        for img_ref in slide_data.images:
            bounds = (content_left, content_top, content_width, 5.0)
            self._insert_image(slide, img_ref.image_name, img_ref.position, bounds)

    def _create_two_column_slide(
        self, prs: Presentation, slide_data: SlideData, template: TemplateConfig
    ) -> None:
        slide = self._add_blank_slide(prs)
        self._set_slide_background(slide, template.background_color)

        self._add_text_box(
            slide,
            TITLE_MARGIN,
            CONTENT_MARGIN,
            11.5,
            0.8,
            slide_data.title,
            template.heading_font,
            template.heading_size,
            template.heading_color,
            bold=True,
        )

        col_top = CONTENT_MARGIN + 1.2
        col_width = 5.5

        self._add_bullet_points(
            slide, TITLE_MARGIN, col_top, col_width, 5.0, slide_data.content, template
        )

        for img_ref in slide_data.images:
            position = img_ref.position if img_ref.position != "center" else "right"
            self._insert_image(slide, img_ref.image_name, position)

        if not slide_data.images and slide_data.right_section:
            self._add_bullet_points(
                slide,
                TITLE_MARGIN + col_width + 0.5,
                col_top,
                col_width,
                5.0,
                slide_data.right_section,
                template,
            )

    def _create_image_focus_slide(
        self, prs: Presentation, slide_data: SlideData, template: TemplateConfig
    ) -> None:
        slide = self._add_blank_slide(prs)
        self._set_slide_background(slide, template.background_color)

        self._add_text_box(
            slide,
            TITLE_MARGIN,
            CONTENT_MARGIN,
            11.5,
            0.8,
            slide_data.title,
            template.heading_font,
            template.heading_size,
            template.heading_color,
            bold=True,
        )

        if slide_data.images:
            self._insert_image(slide, slide_data.images[0].image_name, "center")
        elif slide_data.content:
            self._insert_image(slide, "", "center")

        if slide_data.content:
            self._add_bullet_points(
                slide,
                TITLE_MARGIN,
                5.5,
                11.5,
                1.5,
                slide_data.content[:3],
                template,
            )

    def _create_comparison_slide(
        self, prs: Presentation, slide_data: SlideData, template: TemplateConfig
    ) -> None:
        slide = self._add_blank_slide(prs)
        self._set_slide_background(slide, template.background_color)

        self._add_text_box(
            slide,
            TITLE_MARGIN,
            CONTENT_MARGIN,
            11.5,
            0.8,
            slide_data.title,
            template.heading_font,
            template.heading_size,
            template.heading_color,
            bold=True,
        )

        col_top = CONTENT_MARGIN + 1.2
        col_width = 5.5
        left_items = slide_data.left_section or slide_data.content[: len(slide_data.content) // 2]
        right_items = slide_data.right_section or slide_data.content[len(slide_data.content) // 2 :]

        self._add_text_box(
            slide, TITLE_MARGIN, col_top - 0.5, col_width, 0.4, "Left",
            template.heading_font, template.content_size, template.accent_color, bold=True,
        )
        self._add_bullet_points(slide, TITLE_MARGIN, col_top, col_width, 4.5, left_items, template)

        self._add_text_box(
            slide, TITLE_MARGIN + col_width + 0.5, col_top - 0.5, col_width, 0.4, "Right",
            template.heading_font, template.content_size, template.accent_color, bold=True,
        )
        self._add_bullet_points(
            slide, TITLE_MARGIN + col_width + 0.5, col_top, col_width, 4.5, right_items, template
        )

    def _create_timeline_slide(
        self, prs: Presentation, slide_data: SlideData, template: TemplateConfig
    ) -> None:
        slide = self._add_blank_slide(prs)
        self._set_slide_background(slide, template.background_color)

        self._add_text_box(
            slide,
            TITLE_MARGIN,
            CONTENT_MARGIN,
            11.5,
            0.8,
            slide_data.title,
            template.heading_font,
            template.heading_size,
            template.heading_color,
            bold=True,
        )

        items = slide_data.timeline_items or slide_data.content
        timeline_text = []
        for i, item in enumerate(items, 1):
            timeline_text.append(f"Step {i}: {item}")

        self._add_bullet_points(
            slide, TITLE_MARGIN, CONTENT_MARGIN + 1.2, 11.5, 5.0, timeline_text, template
        )

    def _create_conclusion_slide(
        self, prs: Presentation, slide_data: SlideData, template: TemplateConfig
    ) -> None:
        slide = self._add_blank_slide(prs)
        self._set_slide_background(slide, template.background_color)

        self._add_text_box(
            slide,
            TITLE_MARGIN,
            CONTENT_MARGIN,
            11.5,
            0.8,
            slide_data.title or "Conclusion",
            template.heading_font,
            template.heading_size,
            template.heading_color,
            bold=True,
        )

        self._add_bullet_points(
            slide,
            TITLE_MARGIN,
            CONTENT_MARGIN + 1.2,
            11.5,
            5.0,
            slide_data.content,
            template,
        )

        for img_ref in slide_data.images:
            self._insert_image(slide, img_ref.image_name, img_ref.position)

    def _create_thank_you_slide(
        self, prs: Presentation, slide_data: SlideData, template: TemplateConfig
    ) -> None:
        slide = self._add_blank_slide(prs)
        self._set_slide_background(slide, template.background_color)

        message = slide_data.title or "Thank You"
        subtitle = slide_data.subtitle or slide_data.content[0] if slide_data.content else "Questions?"

        self._add_text_box(
            slide,
            TITLE_MARGIN,
            2.5,
            13.333 - 2 * TITLE_MARGIN,
            1.5,
            message,
            template.title_font,
            template.title_size,
            template.title_color,
            bold=True,
            alignment=PP_ALIGN.CENTER,
        )

        self._add_text_box(
            slide,
            TITLE_MARGIN,
            4.2,
            13.333 - 2 * TITLE_MARGIN,
            1.0,
            subtitle,
            template.body_font,
            template.subtitle_size,
            template.accent_color,
            alignment=PP_ALIGN.CENTER,
        )
